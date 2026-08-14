"""
integrations/presentations.py

Extracts text from an uploaded PowerPoint deck (.pptx) — investment decks,
earnings presentations — so it can be injected into the conversation as
context. Same "raw text IS the grounding" approach as PDFs/spreadsheets,
no parsing/embedding layer.
"""

from __future__ import annotations

import io

from pptx import Presentation

MAX_PPTX_CHARS = 20_000  # matches the cap used for PDFs

SUPPORTED_MIME_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def extract_pptx_text(file_bytes: bytes) -> str | None:
    try:
        prs = Presentation(io.BytesIO(file_bytes))
    except Exception:
        return None

    slides_text = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append(" | ".join(cell.text.strip() for cell in row.cells))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            parts.append(f"[Speaker notes] {slide.notes_slide.notes_text_frame.text.strip()}")
        if parts:
            slides_text.append(f"--- Slide {i} ---\n" + "\n".join(parts))

    text = "\n\n".join(slides_text)
    if not text.strip():
        return None

    if len(text) > MAX_PPTX_CHARS:
        text = text[:MAX_PPTX_CHARS]

    return text
